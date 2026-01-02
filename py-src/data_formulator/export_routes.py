# Export routes for PowerPoint generation
# This endpoint accepts an uploaded image (PNG) and inserts it into a PPTX template
# Place any company templates in a `ppt_templates` folder at the project root or next to this file

import io
import os
import time
from pathlib import Path
from flask import Blueprint, request, send_file, jsonify, current_app
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image

export_bp = Blueprint('export_bp', __name__)

# Default locations to look for templates (module / py-src / repo root / working dir)
MODULE_ROOT = Path(__file__).parent.absolute()
PROJECT_ROOT = MODULE_ROOT.parent
REPO_ROOT = PROJECT_ROOT.parent
PPT_TEMPLATE_DIRS = [
    MODULE_ROOT / 'ppt_templates',
    PROJECT_ROOT / 'ppt_templates',
    REPO_ROOT / 'ppt_templates',
    PROJECT_ROOT / 'templates',
    REPO_ROOT / 'templates',
    PROJECT_ROOT,
    REPO_ROOT,
    Path.cwd(),
]


@export_bp.route('/api/export/list-templates', methods=['GET'])
def list_templates():
    """List all available PowerPoint templates from template directories."""
    templates = set()
    
    for d in PPT_TEMPLATE_DIRS:
        if d.exists() and d.is_dir():
            for file in d.glob('*.pptx'):
                templates.add(file.name)
        elif d.exists() and d.is_file() and d.suffix == '.pptx':
            templates.add(d.name)
    
    return jsonify({
        'status': 'success',
        'templates': sorted(list(templates))
    })


@export_bp.route('/api/export/pptx', methods=['POST'])
def export_pptx():
    # Expect multipart/form-data with 'image' and optional 'template' and 'title'
    if 'image' not in request.files:
        return jsonify({'error': 'No image uploaded'}), 400

    image_file = request.files['image']
    template_name = request.form.get('template', '')
    title = request.form.get('title', '')

    try:
        img_bytes = image_file.read()
        img_io = io.BytesIO(img_bytes)

        prs = None
        if template_name:
            # Find the template in possible template dirs
            current_app.logger.info(f"Looking for template: {template_name}")
            found = None
            for d in PPT_TEMPLATE_DIRS:
                candidate = d / template_name
                current_app.logger.info(f"Checking: {candidate}")
                if candidate.exists():
                    current_app.logger.info(f"Template found at: {candidate}")

                    found = candidate
                    break

            if found:
                current_app.logger.info(f"Loading presentation from: {found}")
                prs = Presentation(str(found))
                current_app.logger.info(f"Template loaded successfully")
            else:
                current_app.logger.error(f"Template '{template_name}' not found in any directory")
                return jsonify({'error': f"Template '{template_name}' not found. Put it in 'ppt_templates' folder."}), 400
        else:
            # create a basic blank presentation
            prs = Presentation()

        # Choose a blank layout for new slide (fallback to 0)
        layout_idx = 6 if len(prs.slide_layouts) > 6 else 0
        blank_layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(blank_layout)

        # Insert image and scale to fit slide while preserving aspect ratio
        img_io.seek(0)
        img = Image.open(img_io)
        width_px, height_px = img.size

        # Slide size in inches
        EMU_PER_INCH = 914400
        slide_width_in = prs.slide_width / EMU_PER_INCH
        slide_height_in = prs.slide_height / EMU_PER_INCH

        # Assume 96 DPI to convert px -> inches
        dpi = 96.0
        img_w_in = width_px / dpi
        img_h_in = height_px / dpi

        # Allow 0.5 inch margin on all sides
        max_w = slide_width_in - 1.0
        max_h = slide_height_in - 1.0
        scale = min(max_w / img_w_in if img_w_in > 0 else 1, max_h / img_h_in if img_h_in > 0 else 1, 1.0)
        final_w = img_w_in * scale
        final_h = img_h_in * scale

        left = (slide_width_in - final_w) / 2
        top = (slide_height_in - final_h) / 2

        # Reset pointer and add picture
        img_io.seek(0)
        slide.shapes.add_picture(img_io, Inches(left), Inches(top), width=Inches(final_w), height=Inches(final_h))

        # Add title textbox if provided
        if title:
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(slide_width_in - 1.0), Inches(0.6))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = title[:200]
            p.font.size = Pt(18)

        out = io.BytesIO()
        prs.save(out)
        out.seek(0)

        filename = f"report-{int(time.time())}.pptx"
        return send_file(out, download_name=filename, as_attachment=True, mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')

    except Exception as e:
        current_app.logger.exception('Failed to create PPTX')
        return jsonify({'error': 'Failed to create PPTX', 'details': str(e)}), 500
